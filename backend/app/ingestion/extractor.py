"""
Document extraction using unstructured (local) and pdfplumber for tables.
Handles PDF, DOCX, PPTX, XLSX, and TXT files — all processing is local.
"""
from __future__ import annotations

import io
import sys
from dataclasses import dataclass, field
from pathlib import Path

# pdfplumber can hit recursion limits on complex PDFs with nested crops
sys.setrecursionlimit(3000)

from app.core.logging import get_logger

logger = get_logger(__name__)


@dataclass
class ExtractedElement:
    """A single extracted element from a document."""
    text: str
    element_type: str  # "narrative", "title", "table", "list_item", "header", "footer"
    page_number: int | None = None
    metadata: dict = field(default_factory=dict)


@dataclass
class ExtractionResult:
    """Result of document extraction."""
    elements: list[ExtractedElement]
    page_count: int
    tables: list[dict]  # Structured table data
    filename: str = ""


def extract_document(file_bytes: bytes, filename: str) -> ExtractionResult:
    """
    Extract text and tables from a document.
    Dispatches to the appropriate extractor based on file extension.
    """
    ext = Path(filename).suffix.lower()
    logger.info("extraction_start", filename=filename, extension=ext, size=len(file_bytes))

    if ext == ".pdf":
        return _extract_pdf(file_bytes, filename)
    elif ext == ".docx":
        return _extract_docx(file_bytes, filename)
    elif ext == ".pptx":
        return _extract_pptx(file_bytes, filename)
    elif ext == ".xlsx":
        return _extract_xlsx(file_bytes, filename)
    elif ext == ".txt":
        return _extract_txt(file_bytes, filename)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def _extract_pdf(file_bytes: bytes, filename: str) -> ExtractionResult:
    """Extract from PDF using camelot for tables, pdfplumber for text, and Gemini for OCR."""
    import tempfile
    import pdfplumber
    import camelot
    from pdf2image import convert_from_bytes
    from google import genai
    from app.core.config import settings

    elements: list[ExtractedElement] = []
    tables: list[dict] = []

    # Initialize Gemini client for OCR fallback on scanned pages
    gemini_client = None
    if settings.GEMINI_API_KEY:
        gemini_client = genai.Client(api_key=settings.GEMINI_API_KEY)

    # camelot requires a file path, so we write the bytes to a temp file
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=True) as temp_pdf:
        temp_pdf.write(file_bytes)
        temp_pdf.flush()

        with pdfplumber.open(temp_pdf.name) as pdf:
            page_count = len(pdf.pages)

            for page_num, page in enumerate(pdf.pages, 1):
                # Check for native text
                raw_text = None
                try:
                    raw_text = page.extract_text(x_tolerance=2, y_tolerance=2)
                except Exception as e:
                    logger.warning("pdfplumber_text_extraction_failed", page=page_num, error=str(e))

                # Threshold: if less than 50 characters, we consider it a scanned/image page
                if raw_text and len(raw_text.strip()) > 50:
                    # ── NATIVE TEXT PAGE ──
                    
                    # 1. Extract tables using camelot
                    try:
                        # Use lattice (border-based) by default, fallback to stream (whitespace-based)
                        camelot_tables = camelot.read_pdf(temp_pdf.name, pages=str(page_num), flavor="lattice")
                        if not camelot_tables:
                            camelot_tables = camelot.read_pdf(temp_pdf.name, pages=str(page_num), flavor="stream")
                    except Exception as e:
                        logger.warning("camelot_extraction_failed", page=page_num, error=str(e))
                        camelot_tables = []

                    for t_idx, table in enumerate(camelot_tables):
                        df = table.df
                        if df.empty or len(df) < 2:
                            continue

                        headers = [str(h).strip() for h in df.iloc[0].tolist()]
                        rows = []
                        for _, row in df.iloc[1:].iterrows():
                            rows.append({
                                headers[j]: str(cell).strip()
                                for j, cell in enumerate(row)
                                if j < len(headers)
                            })

                        tables.append({
                            "page_number": page_num,
                            "table_index": t_idx,
                            "headers": headers,
                            "rows": rows,
                        })

                        # Build a text representation of the table
                        try:
                            table_text = df.to_csv(index=False, sep="|")
                            if table_text.strip():
                                elements.append(ExtractedElement(
                                    text=table_text.strip(),
                                    element_type="table",
                                    page_number=page_num,
                                ))
                        except Exception:
                            pass

                    # 2. Extract narrative text using pdfplumber
                    # We rely on lines to classify text (we don't perfectly crop table bboxes here 
                    # as camelot handles table logic independently)
                    for line in raw_text.split("\n"):
                        line = line.strip()
                        if not line:
                            continue
                        
                        # Heuristic classification
                        if len(line) < 80 and (line.isupper() or line.endswith(":")):
                            etype = "title"
                        elif line.startswith(("•", "-", "*", "–", "◦")):
                            etype = "list_item"
                        else:
                            etype = "narrative"
                            
                        elements.append(ExtractedElement(
                            text=line,
                            element_type=etype,
                            page_number=page_num,
                        ))

                else:
                    # ── SCANNED PAGE (OCR FALLBACK) ──
                    if not gemini_client:
                        logger.warning("gemini_client_missing", msg="Skipping OCR because GEMINI_API_KEY is not set.")
                        continue
                        
                    logger.info("gemini_ocr_fallback", page=page_num, filename=filename)
                    try:
                        # Convert specific page to image
                        images = convert_from_bytes(file_bytes, first_page=page_num, last_page=page_num)
                        if images:
                            image = images[0]
                            prompt = (
                                "Extract all the text from this document page. "
                                "If there are tables, format them strictly as Markdown tables. "
                                "If there are lists, format them as Markdown lists. "
                                "Do not include any introductory or conversational text, just the extracted content."
                            )
                            # Using 1.5 flash as it is highly efficient and budget-friendly for OCR
                            response = gemini_client.models.generate_content(
                                model="gemini-1.5-flash",
                                contents=[prompt, image]
                            )
                            ocr_text = response.text if response.text else ""

                            # Naive parse of Gemini markdown output
                            for block in ocr_text.split("\n\n"):
                                block = block.strip()
                                if not block:
                                    continue
                                if "|" in block and "-|-" in block:
                                    elements.append(ExtractedElement(
                                        text=block,
                                        element_type="table",
                                        page_number=page_num,
                                    ))
                                elif block.startswith(("#", "•", "-", "*")):
                                    elements.append(ExtractedElement(
                                        text=block,
                                        element_type="list_item" if block.startswith(("•", "-", "*")) else "title",
                                        page_number=page_num,
                                    ))
                                else:
                                    elements.append(ExtractedElement(
                                        text=block,
                                        element_type="narrative",
                                        page_number=page_num,
                                    ))
                    except Exception as e:
                        logger.error("gemini_ocr_failed", page=page_num, error=str(e))

    logger.info(
        "extraction_complete",
        filename=filename,
        elements=len(elements),
        tables=len(tables),
        pages=page_count,
    )
    return ExtractionResult(
        elements=elements, page_count=page_count, tables=tables, filename=filename
    )


def _extract_docx(file_bytes: bytes, filename: str) -> ExtractionResult:
    """Extract from DOCX using python-docx."""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    elements: list[ExtractedElement] = []
    tables: list[dict] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style and para.style.name and "heading" in para.style.name.lower():
            etype = "title"
        else:
            etype = "narrative"
        elements.append(ExtractedElement(text=text, element_type=etype))

    for t_idx, table in enumerate(doc.tables):
        if len(table.rows) > 1:
            headers = [cell.text.strip() for cell in table.rows[0].cells]
            rows = []
            for row in table.rows[1:]:
                rows.append({
                    headers[j]: cell.text.strip()
                    for j, cell in enumerate(row.cells)
                    if j < len(headers)
                })
            tables.append({
                "page_number": None,
                "table_index": t_idx,
                "headers": headers,
                "rows": rows,
            })

    return ExtractionResult(
        elements=elements, page_count=len(doc.paragraphs) // 30 or 1,
        tables=tables, filename=filename,
    )


def _extract_pptx(file_bytes: bytes, filename: str) -> ExtractionResult:
    """Extract from PPTX using python-pptx."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    elements: list[ExtractedElement] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        elements.append(ExtractedElement(
                            text=text,
                            element_type="narrative",
                            page_number=slide_num,
                        ))

    return ExtractionResult(
        elements=elements, page_count=len(prs.slides),
        tables=[], filename=filename,
    )


def _extract_xlsx(file_bytes: bytes, filename: str) -> ExtractionResult:
    """Extract from XLSX using openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    elements: list[ExtractedElement] = []
    tables: list[dict] = []

    for sheet_idx, sheet in enumerate(wb.worksheets):
        rows_data = list(sheet.iter_rows(values_only=True))
        if not rows_data:
            continue

        headers = [str(h or "").strip() for h in rows_data[0]]
        rows = []
        for row in rows_data[1:]:
            rows.append({
                headers[j]: str(cell or "").strip()
                for j, cell in enumerate(row)
                if j < len(headers)
            })
        tables.append({
            "page_number": sheet_idx + 1,
            "table_index": 0,
            "headers": headers,
            "rows": rows,
            "sheet_name": sheet.title,
        })
        # Also add as text elements
        for row in rows:
            text = " | ".join(f"{k}: {v}" for k, v in row.items() if v)
            if text:
                elements.append(ExtractedElement(
                    text=text, element_type="table", page_number=sheet_idx + 1
                ))

    wb.close()
    return ExtractionResult(
        elements=elements, page_count=len(wb.worksheets),
        tables=tables, filename=filename,
    )


def _extract_txt(file_bytes: bytes, filename: str) -> ExtractionResult:
    """Extract from plain text."""
    text = file_bytes.decode("utf-8", errors="replace")
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    elements = [
        ExtractedElement(text=p, element_type="narrative", page_number=1)
        for p in paragraphs
    ]
    return ExtractionResult(
        elements=elements, page_count=1, tables=[], filename=filename,
    )
